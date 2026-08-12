from __future__ import annotations

import io
import mimetypes
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pypdf.errors import PdfReadError
from pptx import Presentation
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename

from ..errors import AppError


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".png", ".jpg", ".jpeg"}
INLINE_MIME_TYPES = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
}
MAX_EXTRACTED_CHARS = 400_000
MAX_OFFICE_UNCOMPRESSED_BYTES = 120 * 1024 * 1024


@dataclass(slots=True)
class SourceDocument:
    label: str
    filename: str
    mime_type: str
    text: str = ""
    inline_bytes: bytes | None = None


def _safe_display_name(filename: str) -> str:
    cleaned = secure_filename(filename) or "anexo"
    return cleaned[:120]


def _guard_office_zip(data: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            if len(archive.infolist()) > 5000:
                raise AppError("Um arquivo Office contém itens internos demais para processamento seguro.")
            total = sum(item.file_size for item in archive.infolist())
            if total > MAX_OFFICE_UNCOMPRESSED_BYTES:
                raise AppError("Um arquivo Office se expande além do limite seguro.")
    except zipfile.BadZipFile as exc:
        raise AppError("Um arquivo Office está corrompido ou não é válido.") from exc


def _truncate(text: str) -> str:
    text = re.sub(r"\x00", "", text).strip()
    if len(text) <= MAX_EXTRACTED_CHARS:
        return text
    return text[:MAX_EXTRACTED_CHARS] + "\n\n[CONTEÚDO TRUNCADO PELO LIMITE TÉCNICO DO APLICATIVO]"


def _extract_docx(data: bytes) -> str:
    _guard_office_zip(data)
    doc = Document(io.BytesIO(data))
    blocks: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            blocks.append(paragraph.text.strip())
    for table_index, table in enumerate(doc.tables, start=1):
        blocks.append(f"\n[TABELA {table_index}]")
        for row in table.rows:
            seen: set[int] = set()
            values: list[str] = []
            for cell in row.cells:
                cell_id = id(cell._tc)
                if cell_id in seen:
                    continue
                seen.add(cell_id)
                value = " ".join(p.text.strip() for p in cell.paragraphs if p.text.strip())
                if value:
                    values.append(value)
            if values:
                blocks.append(" | ".join(values))
    return _truncate("\n".join(blocks))


def _extract_xlsx(data: bytes) -> str:
    _guard_office_zip(data)
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[str] = []
    extracted_chars = 0
    try:
        for sheet in workbook.worksheets:
            heading = f"\n[PLANILHA: {sheet.title}]"
            blocks.append(heading)
            extracted_chars += len(heading)
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    row_text = " | ".join(values)
                    blocks.append(row_text)
                    extracted_chars += len(row_text)
                if extracted_chars > MAX_EXTRACTED_CHARS:
                    return _truncate("\n".join(blocks))
    finally:
        workbook.close()
    return _truncate("\n".join(blocks))


def _extract_pptx(data: bytes) -> str:
    _guard_office_zip(data)
    presentation = Presentation(io.BytesIO(data))
    blocks: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks.append(f"\n[SLIDE {slide_number}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                blocks.append(text)
    return _truncate("\n".join(blocks))


def _extract_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return _truncate(data.decode(encoding))
        except UnicodeDecodeError:
            continue
    raise AppError("Não foi possível ler um arquivo de texto.")


def _extract_pdf(data: bytes) -> str:
    """Extrai PDFs textuais em memória; PDFs digitalizados permanecem inline."""

    try:
        reader = PdfReader(io.BytesIO(data), strict=False)
        if reader.is_encrypted:
            try:
                unlocked = reader.decrypt("")
            except Exception as exc:
                raise AppError("O PDF está protegido por senha e não pode ser processado.") from exc
            if not unlocked:
                raise AppError("O PDF está protegido por senha e não pode ser processado.")

        blocks: list[str] = []
        extracted_chars = 0
        for page_number, page in enumerate(reader.pages[:120], start=1):
            text = (page.extract_text() or "").strip()
            if text:
                page_text = f"\n[PÁGINA {page_number}]\n{text}"
                blocks.append(page_text)
                extracted_chars += len(page_text)
            if extracted_chars >= MAX_EXTRACTED_CHARS:
                break
        return _truncate("\n".join(blocks)) if blocks else ""
    except AppError:
        raise
    except (PdfReadError, ValueError, OSError) as exc:
        raise AppError("Um arquivo PDF está corrompido ou não é válido.") from exc


def process_upload(
    storage: FileStorage,
    label: str,
    *,
    max_file_bytes: int,
    prefer_pdf_text: bool = False,
) -> SourceDocument:
    original_name = storage.filename or ""
    extension = Path(original_name).suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_EXTENSIONS))
        raise AppError(f"Formato não permitido em '{original_name}'. Use: {allowed}.")

    data = storage.read(max_file_bytes + 1)
    if not data:
        raise AppError(f"O anexo '{original_name}' está vazio.")
    if len(data) > max_file_bytes:
        max_mb = max_file_bytes // (1024 * 1024)
        raise AppError(f"O anexo '{original_name}' ultrapassa {max_mb} MB.")

    filename = _safe_display_name(original_name)
    mime_type = INLINE_MIME_TYPES.get(extension) or storage.mimetype or mimetypes.guess_type(filename)[0] or "text/plain"

    if extension in INLINE_MIME_TYPES:
        if extension == ".pdf" and not data.startswith(b"%PDF"):
            raise AppError(f"O anexo '{original_name}' não contém um PDF válido.")
        if extension == ".png" and not data.startswith(b"\x89PNG\r\n\x1a\n"):
            raise AppError(f"O anexo '{original_name}' não contém uma imagem PNG válida.")
        if extension in {".jpg", ".jpeg"} and not data.startswith(b"\xff\xd8\xff"):
            raise AppError(f"O anexo '{original_name}' não contém uma imagem JPEG válida.")
        if extension == ".pdf" and prefer_pdf_text:
            extracted = _extract_pdf(data)
            if extracted:
                return SourceDocument(label=label, filename=filename, mime_type="text/plain", text=extracted)
        return SourceDocument(label=label, filename=filename, mime_type=mime_type, inline_bytes=data)
    if extension == ".docx":
        text = _extract_docx(data)
    elif extension == ".xlsx":
        text = _extract_xlsx(data)
    elif extension == ".pptx":
        text = _extract_pptx(data)
    else:
        text = _extract_text(data)

    if not text.strip():
        raise AppError(f"Não foi encontrado texto legível em '{original_name}'.")
    return SourceDocument(label=label, filename=filename, mime_type="text/plain", text=text)


def collect_sources(
    upload_groups: list[tuple[str, list[FileStorage]]],
    *,
    max_file_bytes: int,
    max_files: int,
    prefer_pdf_text: bool = False,
) -> list[SourceDocument]:
    flattened = [(label, item) for label, items in upload_groups for item in items if item and item.filename]
    if len(flattened) > max_files:
        raise AppError(f"Envie no máximo {max_files} anexos por geração.")
    return [
        process_upload(item, label, max_file_bytes=max_file_bytes, prefer_pdf_text=prefer_pdf_text)
        for label, item in flattened
    ]
